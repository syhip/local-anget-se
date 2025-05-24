"""
文档转换模块 - 混合方案实现

该模块实现了一个混合文档转换引擎，主要使用MarkItDown，
备选使用Marker和专用Python库，以实现高质量的文档转换。

支持格式:
- Word (.doc, .docx)
- Excel (.xls, .xlsx)
- PowerPoint (.ppt, .pptx)
- PDF (.pdf)
- 图像 (.jpg, .png, .tiff)
- HTML (.html, .htm)
- 文本 (.txt, .csv, .json, .xml)

作者: Manus AI
日期: 2025-05-23
"""

import os
import sys
import logging
import tempfile
import shutil
from pathlib import Path
from enum import Enum
from typing import Optional, List, Dict, Any, Union, Tuple

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

class ConversionEngine(Enum):
    """文档转换引擎枚举"""
    MARKITDOWN = "markitdown"
    MARKER = "marker"
    PYTHON_LIBS = "python_libs"
    AUTO = "auto"  # 自动选择最佳引擎

class DocumentType(Enum):
    """文档类型枚举"""
    WORD = "word"
    EXCEL = "excel"
    POWERPOINT = "powerpoint"
    PDF = "pdf"
    IMAGE = "image"
    HTML = "html"
    TEXT = "text"
    UNKNOWN = "unknown"

class ConversionResult:
    """转换结果类"""
    def __init__(
        self, 
        success: bool, 
        markdown_content: str = "", 
        engine_used: str = "",
        error_message: str = "",
        metadata: Dict[str, Any] = None,
        images: List[Dict[str, Any]] = None
    ):
        self.success = success
        self.markdown_content = markdown_content
        self.engine_used = engine_used
        self.error_message = error_message
        self.metadata = metadata or {}
        self.images = images or []
    
    def __str__(self) -> str:
        if self.success:
            return f"转换成功 [引擎: {self.engine_used}] - 内容长度: {len(self.markdown_content)}字符"
        else:
            return f"转换失败 [引擎: {self.engine_used}] - 错误: {self.error_message}"

class DocumentConverter:
    """文档转换器主类"""
    
    def __init__(self, config: Dict[str, Any] = None):
        """
        初始化文档转换器
        
        Args:
            config: 配置字典，可包含以下键:
                - default_engine: 默认转换引擎
                - temp_dir: 临时文件目录
                - markitdown_options: MarkItDown特定选项
                - marker_options: Marker特定选项
                - python_libs_options: Python库特定选项
        """
        self.config = config or {}
        self.default_engine = ConversionEngine(self.config.get("default_engine", "auto"))
        self.temp_dir = self.config.get("temp_dir", tempfile.gettempdir())
        
        # 确保临时目录存在
        os.makedirs(self.temp_dir, exist_ok=True)
        
        # 初始化各引擎的配置
        self.markitdown_options = self.config.get("markitdown_options", {})
        self.marker_options = self.config.get("marker_options", {})
        self.python_libs_options = self.config.get("python_libs_options", {})
        
        # 检查依赖
        self._check_dependencies()
    
    def _check_dependencies(self) -> None:
        """检查必要的依赖是否已安装"""
        try:
            # 尝试导入必要的库
            import importlib
            
            # 检查MarkItDown
            try:
                importlib.import_module("markitdown")
                self.markitdown_available = True
                logger.info("MarkItDown已安装")
            except ImportError:
                self.markitdown_available = False
                logger.warning("MarkItDown未安装，将无法使用该引擎")
            
            # 检查Marker
            try:
                importlib.import_module("marker")
                self.marker_available = True
                logger.info("Marker已安装")
            except ImportError:
                self.marker_available = False
                logger.warning("Marker未安装，将无法使用该引擎")
            
            # 检查Python库
            python_libs = {
                "python-docx": "docx",
                "pandas": "pandas",
                "PyPDF2": "PyPDF2",
                "pdfminer.six": "pdfminer"
            }
            
            self.python_libs_available = {}
            for lib_name, module_name in python_libs.items():
                try:
                    importlib.import_module(module_name)
                    self.python_libs_available[lib_name] = True
                    logger.info(f"{lib_name}已安装")
                except ImportError:
                    self.python_libs_available[lib_name] = False
                    logger.warning(f"{lib_name}未安装，某些功能可能受限")
        
        except Exception as e:
            logger.error(f"检查依赖时出错: {str(e)}")
    
    def _detect_document_type(self, file_path: str) -> DocumentType:
        """
        根据文件扩展名检测文档类型
        
        Args:
            file_path: 文件路径
            
        Returns:
            DocumentType: 文档类型枚举
        """
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext in ['.doc', '.docx']:
            return DocumentType.WORD
        elif ext in ['.xls', '.xlsx']:
            return DocumentType.EXCEL
        elif ext in ['.ppt', '.pptx']:
            return DocumentType.POWERPOINT
        elif ext == '.pdf':
            return DocumentType.PDF
        elif ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff']:
            return DocumentType.IMAGE
        elif ext in ['.html', '.htm']:
            return DocumentType.HTML
        elif ext in ['.txt', '.csv', '.json', '.xml']:
            return DocumentType.TEXT
        else:
            return DocumentType.UNKNOWN
    
    def _select_best_engine(self, file_path: str) -> ConversionEngine:
        """
        根据文档类型自动选择最佳转换引擎
        
        Args:
            file_path: 文件路径
            
        Returns:
            ConversionEngine: 选择的转换引擎
        """
        doc_type = self._detect_document_type(file_path)
        
        # 如果MarkItDown可用，优先使用
        if self.markitdown_available:
            return ConversionEngine.MARKITDOWN
        
        # 如果是PDF或图像，且Marker可用，使用Marker
        if doc_type in [DocumentType.PDF, DocumentType.IMAGE] and self.marker_available:
            return ConversionEngine.MARKER
        
        # 其他情况使用Python库
        return ConversionEngine.PYTHON_LIBS
    
    def _convert_with_markitdown(self, file_path: str) -> ConversionResult:
        """
        使用MarkItDown转换文档
        
        Args:
            file_path: 文件路径
            
        Returns:
            ConversionResult: 转换结果
        """
        try:
            from markitdown import MarkItDown
            
            # 创建MarkItDown实例
            md = MarkItDown(enable_plugins=self.markitdown_options.get("enable_plugins", False))
            
            # 执行转换
            result = md.convert(file_path)
            
            # 返回结果
            return ConversionResult(
                success=True,
                markdown_content=result.text_content,
                engine_used="MarkItDown",
                metadata=result.metadata if hasattr(result, 'metadata') else {},
                images=result.images if hasattr(result, 'images') else []
            )
        
        except Exception as e:
            logger.error(f"使用MarkItDown转换时出错: {str(e)}")
            return ConversionResult(
                success=False,
                engine_used="MarkItDown",
                error_message=str(e)
            )
    
    def _convert_with_marker(self, file_path: str) -> ConversionResult:
        """
        使用Marker转换文档
        
        Args:
            file_path: 文件路径
            
        Returns:
            ConversionResult: 转换结果
        """
        try:
            from marker.converters.pdf import PdfConverter
            from marker.models import create_model_dict
            from marker.output import text_from_rendered
            
            # 创建转换器
            converter = PdfConverter(
                artifact_dict=create_model_dict(),
                **self.marker_options
            )
            
            # 执行转换
            rendered = converter(file_path)
            text, metadata, images = text_from_rendered(rendered)
            
            # 返回结果
            return ConversionResult(
                success=True,
                markdown_content=text,
                engine_used="Marker",
                metadata=metadata or {},
                images=images or []
            )
        
        except Exception as e:
            logger.error(f"使用Marker转换时出错: {str(e)}")
            return ConversionResult(
                success=False,
                engine_used="Marker",
                error_message=str(e)
            )
    
    def _convert_with_python_libs(self, file_path: str) -> ConversionResult:
        """
        使用Python库转换文档
        
        Args:
            file_path: 文件路径
            
        Returns:
            ConversionResult: 转换结果
        """
        doc_type = self._detect_document_type(file_path)
        
        try:
            if doc_type == DocumentType.WORD:
                return self._convert_word(file_path)
            elif doc_type == DocumentType.EXCEL:
                return self._convert_excel(file_path)
            elif doc_type == DocumentType.POWERPOINT:
                return self._convert_powerpoint(file_path)
            elif doc_type == DocumentType.PDF:
                return self._convert_pdf(file_path)
            elif doc_type == DocumentType.HTML:
                return self._convert_html(file_path)
            elif doc_type == DocumentType.TEXT:
                return self._convert_text(file_path)
            else:
                return ConversionResult(
                    success=False,
                    engine_used="Python库",
                    error_message=f"不支持的文档类型: {doc_type}"
                )
        
        except Exception as e:
            logger.error(f"使用Python库转换时出错: {str(e)}")
            return ConversionResult(
                success=False,
                engine_used="Python库",
                error_message=str(e)
            )
    
    def _convert_word(self, file_path: str) -> ConversionResult:
        """使用python-docx转换Word文档"""
        try:
            import docx
            
            doc = docx.Document(file_path)
            markdown_content = []
            
            # 处理段落
            for para in doc.paragraphs:
                if para.style.name.startswith('Heading'):
                    level = int(para.style.name[-1])
                    markdown_content.append(f"{'#' * level} {para.text}\n")
                else:
                    markdown_content.append(f"{para.text}\n\n")
            
            # 处理表格
            for table in doc.tables:
                markdown_table = []
                # 添加表头
                header_row = []
                for cell in table.rows[0].cells:
                    header_row.append(cell.text.strip())
                markdown_table.append("| " + " | ".join(header_row) + " |")
                
                # 添加分隔行
                markdown_table.append("| " + " | ".join(["---"] * len(header_row)) + " |")
                
                # 添加数据行
                for row in table.rows[1:]:
                    data_row = []
                    for cell in row.cells:
                        data_row.append(cell.text.strip())
                    markdown_table.append("| " + " | ".join(data_row) + " |")
                
                markdown_content.append("\n".join(markdown_table) + "\n\n")
            
            return ConversionResult(
                success=True,
                markdown_content="\n".join(markdown_content),
                engine_used="python-docx"
            )
        
        except Exception as e:
            logger.error(f"转换Word文档时出错: {str(e)}")
            return ConversionResult(
                success=False,
                engine_used="python-docx",
                error_message=str(e)
            )
    
    def _convert_excel(self, file_path: str) -> ConversionResult:
        """使用pandas转换Excel文档"""
        try:
            import pandas as pd
            
            # 读取所有工作表
            excel_file = pd.ExcelFile(file_path)
            sheet_names = excel_file.sheet_names
            
            markdown_content = []
            
            # 处理每个工作表
            for sheet_name in sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                # 添加工作表标题
                markdown_content.append(f"## 工作表: {sheet_name}\n")
                
                # 转换为Markdown表格
                markdown_table = df.to_markdown(index=False)
                markdown_content.append(markdown_table + "\n\n")
            
            return ConversionResult(
                success=True,
                markdown_content="\n".join(markdown_content),
                engine_used="pandas"
            )
        
        except Exception as e:
            logger.error(f"转换Excel文档时出错: {str(e)}")
            return ConversionResult(
                success=False,
                engine_used="pandas",
                error_message=str(e)
            )
    
    def _convert_powerpoint(self, file_path: str) -> ConversionResult:
        """转换PowerPoint文档"""
        try:
            import pptx
            
            presentation = pptx.Presentation(file_path)
            markdown_content = []
            
            # 处理每张幻灯片
            for i, slide in enumerate(presentation.slides):
                markdown_content.append(f"## 幻灯片 {i+1}\n")
                
                # 处理幻灯片标题
                if slide.shapes.title:
                    markdown_content.append(f"### {slide.shapes.title.text}\n")
                
                # 处理文本框
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        if shape != slide.shapes.title:  # 避免重复标题
                            markdown_content.append(f"{shape.text}\n\n")
                
                markdown_content.append("\n")
            
            return ConversionResult(
                success=True,
                markdown_content="\n".join(markdown_content),
                engine_used="python-pptx"
            )
        
        except Exception as e:
            logger.error(f"转换PowerPoint文档时出错: {str(e)}")
            return ConversionResult(
                success=False,
                engine_used="python-pptx",
                error_message=str(e)
            )
    
    def _convert_pdf(self, file_path: str) -> ConversionResult:
        """使用PyPDF2和pdfminer转换PDF文档"""
        try:
            from PyPDF2 import PdfReader
            from pdfminer.high_level import extract_text
            
            # 尝试使用pdfminer提取文本
            try:
                text = extract_text(file_path)
                engine = "pdfminer"
            except Exception:
                # 如果pdfminer失败，尝试使用PyPDF2
                reader = PdfReader(file_path)
                text_parts = []
                
                for page in reader.pages:
                    text_parts.append(page.extract_text())
                
                text = "\n\n".join(text_parts)
                engine = "PyPDF2"
            
            # 简单格式化
            lines = text.split("\n")
            markdown_content = []
            
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                
                # 尝试检测标题
                if len(line) < 100 and line.endswith(":"):
                    markdown_content.append(f"### {line}\n")
                else:
                    markdown_content.append(f"{line}\n\n")
            
            return ConversionResult(
                success=True,
                markdown_content="\n".join(markdown_content),
                engine_used=engine
            )
        
        except Exception as e:
            logger.error(f"转换PDF文档时出错: {str(e)}")
            return ConversionResult(
                success=False,
                engine_used="PDF库",
                error_message=str(e)
            )
    
    def _convert_html(self, file_path: str) -> ConversionResult:
        """转换HTML文档"""
        try:
            from bs4 import BeautifulSoup
            import html2text
            
            # 读取HTML文件
            with open(file_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            # 使用html2text转换
            h = html2text.HTML2Text()
            h.ignore_links = False
            h.ignore_images = False
            h.ignore_tables = False
            
            markdown_content = h.handle(html_content)
            
            return ConversionResult(
                success=True,
                markdown_content=markdown_content,
                engine_used="html2text"
            )
        
        except Exception as e:
            logger.error(f"转换HTML文档时出错: {str(e)}")
            return ConversionResult(
                success=False,
                engine_used="html2text",
                error_message=str(e)
            )
    
    def _convert_text(self, file_path: str) -> ConversionResult:
        """转换文本文档"""
        try:
            # 读取文本文件
            with open(file_path, 'r', encoding='utf-8') as f:
                text_content = f.read()
            
            # 对于CSV、JSON、XML等格式，可以进行特殊处理
            ext = os.path.splitext(file_path)[1].lower()
            
            if ext == '.csv':
                import pandas as pd
                df = pd.read_csv(file_path)
                markdown_content = df.to_markdown(index=False)
                engine = "pandas (CSV)"
            
            elif ext == '.json':
                import json
                json_data = json.loads(text_content)
                markdown_content = f"```json\n{json.dumps(json_data, indent=2, ensure_ascii=False)}\n```"
                engine = "json"
            
            elif ext == '.xml':
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(text_content, 'xml')
                markdown_content = f"```xml\n{soup.prettify()}\n```"
                engine = "BeautifulSoup (XML)"
            
            else:
                # 普通文本文件，直接返回内容
                markdown_content = text_content
                engine = "text"
            
            return ConversionResult(
                success=True,
                markdown_content=markdown_content,
                engine_used=engine
            )
        
        except Exception as e:
            logger.error(f"转换文本文档时出错: {str(e)}")
            return ConversionResult(
                success=False,
                engine_used="文本处理",
                error_message=str(e)
            )
    
    def convert(self, file_path: str, engine: ConversionEngine = None) -> ConversionResult:
        """
        转换文档为Markdown格式
        
        Args:
            file_path: 文件路径
            engine: 指定使用的转换引擎，如果为None则使用默认引擎或自动选择
            
        Returns:
            ConversionResult: 转换结果
        """
        # 检查文件是否存在
        if not os.path.exists(file_path):
            return ConversionResult(
                success=False,
                error_message=f"文件不存在: {file_path}"
            )
        
        # 确定使用的引擎
        engine = engine or self.default_engine
        
        # 如果是自动模式，选择最佳引擎
        if engine == ConversionEngine.AUTO:
            engine = self._select_best_engine(file_path)
        
        logger.info(f"使用引擎 {engine.value} 转换文件: {file_path}")
        
        # 根据引擎执行转换
        if engine == ConversionEngine.MARKITDOWN:
            if not self.markitdown_available:
                logger.warning("MarkItDown未安装，尝试使用备选引擎")
                return self.convert(file_path, ConversionEngine.MARKER)
            return self._convert_with_markitdown(file_path)
        
        elif engine == ConversionEngine.MARKER:
            if not self.marker_available:
                logger.warning("Marker未安装，尝试使用备选引擎")
                return self.convert(file_path, ConversionEngine.PYTHON_LIBS)
            return self._convert_with_marker(file_path)
        
        elif engine == ConversionEngine.PYTHON_LIBS:
            return self._convert_with_python_libs(file_path)
        
        else:
            return ConversionResult(
                success=False,
                error_message=f"未知引擎: {engine}"
            )
    
    def batch_convert(self, file_paths: List[str], engine: ConversionEngine = None) -> Dict[str, ConversionResult]:
        """
        批量转换文档
        
        Args:
            file_paths: 文件路径列表
            engine: 指定使用的转换引擎，如果为None则使用默认引擎或自动选择
            
        Returns:
            Dict[str, ConversionResult]: 文件路径到转换结果的映射
        """
        results = {}
        
        for file_path in file_paths:
            results[file_path] = self.convert(file_path, engine)
        
        return results
    
    def convert_directory(
        self, 
        directory: str, 
        output_directory: str = None,
        recursive: bool = False,
        file_extensions: List[str] = None,
        engine: ConversionEngine = None
    ) -> Dict[str, ConversionResult]:
        """
        转换目录中的所有文档
        
        Args:
            directory: 输入目录
            output_directory: 输出目录，如果为None则不保存文件
            recursive: 是否递归处理子目录
            file_extensions: 要处理的文件扩展名列表，如果为None则处理所有支持的格式
            engine: 指定使用的转换引擎，如果为None则使用默认引擎或自动选择
            
        Returns:
            Dict[str, ConversionResult]: 文件路径到转换结果的映射
        """
        # 检查目录是否存在
        if not os.path.isdir(directory):
            logger.error(f"目录不存在: {directory}")
            return {}
        
        # 如果指定了输出目录，确保它存在
        if output_directory:
            os.makedirs(output_directory, exist_ok=True)
        
        # 确定要处理的文件扩展名
        if file_extensions is None:
            file_extensions = [
                '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx', 
                '.pdf', '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff',
                '.html', '.htm', '.txt', '.csv', '.json', '.xml'
            ]
        
        # 查找所有符合条件的文件
        file_paths = []
        
        if recursive:
            for root, _, files in os.walk(directory):
                for file in files:
                    if any(file.lower().endswith(ext) for ext in file_extensions):
                        file_paths.append(os.path.join(root, file))
        else:
            for file in os.listdir(directory):
                file_path = os.path.join(directory, file)
                if os.path.isfile(file_path) and any(file.lower().endswith(ext) for ext in file_extensions):
                    file_paths.append(file_path)
        
        # 批量转换文件
        results = self.batch_convert(file_paths, engine)
        
        # 如果指定了输出目录，保存转换结果
        if output_directory:
            for file_path, result in results.items():
                if result.success:
                    # 计算相对路径，保持目录结构
                    rel_path = os.path.relpath(file_path, directory)
                    output_path = os.path.join(output_directory, os.path.splitext(rel_path)[0] + '.md')
                    
                    # 确保输出目录存在
                    os.makedirs(os.path.dirname(output_path), exist_ok=True)
                    
                    # 保存Markdown内容
                    with open(output_path, 'w', encoding='utf-8') as f:
                        f.write(result.markdown_content)
                    
                    logger.info(f"已保存转换结果到: {output_path}")
        
        return results

# 命令行接口
if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(description="文档转换工具")
    parser.add_argument("input", help="输入文件或目录路径")
    parser.add_argument("-o", "--output", help="输出文件或目录路径")
    parser.add_argument("-e", "--engine", choices=["markitdown", "marker", "python_libs", "auto"], default="auto", help="使用的转换引擎")
    parser.add_argument("-r", "--recursive", action="store_true", help="递归处理子目录")
    parser.add_argument("-f", "--format", choices=["markdown", "json"], default="markdown", help="输出格式")
    parser.add_argument("-v", "--verbose", action="store_true", help="显示详细日志")
    
    args = parser.parse_args()
    
    # 设置日志级别
    if args.verbose:
        logging.getLogger().setLevel(logging.DEBUG)
    
    # 创建转换器
    converter = DocumentConverter(config={
        "default_engine": args.engine
    })
    
    # 执行转换
    if os.path.isdir(args.input):
        # 处理目录
        if not args.output:
            args.output = args.input + "_converted"
        
        results = converter.convert_directory(
            args.input,
            args.output,
            args.recursive
        )
        
        # 打印统计信息
        success_count = sum(1 for result in results.values() if result.success)
        total_count = len(results)
        
        print(f"转换完成: {success_count}/{total_count} 个文件成功")
    
    else:
        # 处理单个文件
        result = converter.convert(args.input)
        
        if result.success:
            if args.output:
                # 保存到文件
                with open(args.output, 'w', encoding='utf-8') as f:
                    f.write(result.markdown_content)
                print(f"转换成功，已保存到: {args.output}")
            else:
                # 输出到控制台
                print(result.markdown_content)
        else:
            print(f"转换失败: {result.error_message}")
            sys.exit(1)
